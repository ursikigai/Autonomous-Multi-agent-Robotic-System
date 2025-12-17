#!/usr/bin/env python3
import os
from pptx import Presentation
from pptx.util import Inches, Pt

BASE = "/Users/ranjeet/thesis_multiagent/results/final_package"
OUT_PPT = os.path.join(BASE, "final_report.pptx")

prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Autonomous Navigation Report"
slide.placeholders[1].text = "SLAM + YOLO + Behavior + Risk\nGenerated Automatically"

# Images
slides = [
    ("Top-Down SLAM Trajectory", "traj_topdown.png"),
    ("3D SLAM Trajectory", "traj_3d.png"),
    ("SLAM + YOLO Fusion Tracks", "fusion_tracks.png"),
    ("Object Boxes Snapshot", "boxes_snapshot.png"),
    ("Near-Miss Summary", "near_miss_events.png"),
    ("Risk-Time Scatter Plot", "risk_plot.png"),
    ("Future Risk Map", "risk_map_future.png")
]

for label, filename in slides:
    path = os.path.join(BASE, filename)
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = label

    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(0.5), Inches(1.2), width=Inches(9))
    else:
        # Add text placeholder instead of image
        tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
        tx.text = "(Missing image)"

prs.save(OUT_PPT)
print("PPT saved to:", OUT_PPT)


